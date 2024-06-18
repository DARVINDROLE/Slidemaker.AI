from pptx import Presentation
from pptx.util import Inches

import openai
openai.api_key = "sk-proj-hWduDtOJade1ddQiyLSVT3BlbkFJSNo87tvni9fBqqejMqTN"

def generate_response(prompt, max_words=100):
    try:
        response = openai.Completion.create(
            model="gpt-3.5-turbo",  # Adjust to the model you're using, like "gpt-3.5-turbo"
            prompt=prompt,
            max_tokens=200,  # Limit to a reasonable number of tokens for initial response
            temperature=0.5,
        )
        generated_text = response.choices[0].text.strip()

        # Limit the response to the first 100 words
        words = generated_text.split()
        limited_response = ' '.join(words[:max_words])

        return limited_response
    except Exception as e:
        # Return more details about the exception for debugging
        return f"Error generating response. Exception: {e}"

# Example usage

topic = str(input("Enter the topic for PPT: "))
prompt = f"make 10 slide presentation on topic {topic}"
presentation_data='''Slide 1: Title Slide

Title: Introduction to Neural Networks
Subtitle: Understanding the Basics
Your Name
Date
Slide 2: Introduction to Neural Networks

Definition of Neural Networks
Brief history and development
Importance in modern technology
Slide 3: Structure of a Neuron

Biological inspiration
Components: dendrites, cell body, axon
Artificial neuron model (Perceptron)
Slide 4: Types of Neural Networks

Feedforward Neural Networks (FNN)
Recurrent Neural Networks (RNN)
Convolutional Neural Networks (CNN)
Applications and differences
Slide 5: Working Principle of Neural Networks

Input layer, hidden layers, output layer
Activation function
Forward propagation
Slide 6: Training Neural Networks

Loss function and optimization
Backpropagation algorithm
Gradient descent
Slide 7: Common Activation Functions

Sigmoid
ReLU (Rectified Linear Unit)
Tanh (Hyperbolic Tangent)
Comparison and usage scenarios
Slide 8: Applications of Neural Networks

Image and speech recognition
Natural language processing
Autonomous vehicles
Healthcare and diagnostics
Slide 9: Challenges and Limitations

Overfitting and underfitting
Computational complexity
Data requirements and biases
Slide 10: Future of Neural Networks

Advances in deep learning
Integration with other technologies (AI, IoT)
Ethical considerations and societal impacts
Slide 11: Conclusion

Recap main points
Importance of neural networks in future technology
Thank you!''' #generate_response(prompt, max_words=100)
slides = presentation_data.strip().split("Slide ")[1:]  # Split by "Slide " and ignore the first empty string

    # Create a presentation object
prs = Presentation()

    # Add slides based on the presentation data
for slide_data in slides:
        slide_content = slide_data.split("\n\n")
        slide_title = slide_content[0]
        slide_text = "\n".join(slide_content[1:])  # Join the content lines

        slide_layout = prs.slide_layouts[1]  # Assuming we use the second slide layout for content slides
        slide = prs.slides.add_slide(slide_layout)

        # Set slide title
        title_placeholder = slide.shapes.title
        if title_placeholder is not None:
           title_placeholder.text = slide_title

        # Set slide content
        content_placeholder = slide.placeholders[1]  # Assuming the second placeholder is for content
        if content_placeholder is not None:
           content_placeholder.text = slide_text

        # Save the presentation to a file
file_path = f"{topic}.pptx"
prs.save(file_path)

    